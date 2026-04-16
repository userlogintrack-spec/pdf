import io
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree


_BULLET_CHARS = set("•●◦▪▫■□◆◇‣⁃※*⋅∙")


def _strip_bullet_prefix(text):
    """If a line starts with a bullet glyph, return (bullet_char, rest)."""
    stripped = text.lstrip()
    if not stripped:
        return None, text
    leading_ws = text[:len(text) - len(stripped)]
    first = stripped[0]
    if first in _BULLET_CHARS:
        # Also strip the spaces after the bullet
        after = stripped[1:].lstrip()
        return first, leading_ws + after
    return None, text


def _rects_for_links(page):
    """Return list of (rect, uri) for URI link annotations on the page."""
    links = []
    for lk in page.get_links():
        if lk.get("kind") == fitz.LINK_URI and lk.get("uri"):
            links.append((lk["from"], lk["uri"]))
    return links


def _find_link_for_span(span_bbox, links):
    """Return URI if span's bbox is inside any link annotation's rect."""
    sx0, sy0, sx1, sy1 = span_bbox
    sxc = (sx0 + sx1) / 2
    syc = (sy0 + sy1) / 2
    for rect, uri in links:
        if rect.x0 <= sxc <= rect.x1 and rect.y0 <= syc <= rect.y1:
            return uri
    return None


def _set_run_hyperlink(run, uri):
    """Attach a clickable hyperlink to a text run."""
    try:
        run.hyperlink.address = uri
    except Exception:
        pass


def _set_run_invisible(run):
    """
    Set a text run to be fully transparent (alpha=0) while keeping it
    selectable, searchable, and editable. This is the same technique used
    by Adobe Acrobat / iLovePDF for "hidden text under image" layers.
    """
    rPr = run._r.get_or_add_rPr()
    # Remove any existing solidFill
    for existing in rPr.findall(qn('a:solidFill')):
        rPr.remove(existing)
    solid_fill = etree.SubElement(rPr, qn('a:solidFill'))
    srgb = etree.SubElement(solid_fill, qn('a:srgbClr'), val='000000')
    etree.SubElement(srgb, qn('a:alpha'), val='0')


def convert_pdf_to_ppt(file_path, mode='hybrid', dpi=250, pages=None):
    """
    Convert PDF to PowerPoint presentation.

    Modes:
      - 'image'    : Each page rendered as a high-res image (pixel-perfect, not editable)
      - 'editable' : Extracts text, images, shapes as native PPT elements (fully editable)
      - 'hybrid'   : Image background for pixel-perfect layout + editable text overlay on top

    Args:
        file_path: Path to the PDF file
        mode: Conversion mode ('image', 'editable', 'hybrid')
        dpi: Resolution for image rendering (72-600)
        pages: Optional list of 0-based page numbers to convert, None = all pages
    """
    pdf_doc = fitz.open(file_path)
    prs = Presentation()

    dpi = max(72, min(600, dpi))
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)

    page_indices = _resolve_pages(pdf_doc.page_count, pages)

    for page_num in page_indices:
        page = pdf_doc[page_num]
        pw = page.rect.width
        ph = page.rect.height

        slide_w, slide_h = _calc_slide_dimensions(pw, ph)
        prs.slide_width = slide_w
        prs.slide_height = slide_h

        # Use blank layout (index 6)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if mode == 'image':
            _build_image_slide(slide, page, mat, slide_w, slide_h)
        elif mode == 'editable':
            _build_editable_slide(slide, page, pdf_doc, slide_w, slide_h, pw, ph)
        else:
            # hybrid: image background + editable text on top
            _build_hybrid_slide(slide, page, pdf_doc, mat, slide_w, slide_h, pw, ph)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    pdf_doc.close()
    return output.getvalue()


# ─── Page resolution ────────────────────────────────────────────────

def _resolve_pages(total, pages):
    """Return sorted list of valid 0-based page indices."""
    if pages is None:
        return list(range(total))
    valid = [p for p in pages if 0 <= p < total]
    return sorted(set(valid)) if valid else list(range(total))


# ─── Slide dimensions ──────────────────────────────────────────────

def _calc_slide_dimensions(pw, ph):
    """Calculate slide dimensions preserving PDF aspect ratio."""
    ratio = pw / ph
    if ratio > 1.2:
        # Landscape - widescreen
        slide_w = Inches(13.333)
        slide_h = int(slide_w / ratio)
    elif ratio < 0.8:
        # Portrait
        slide_h = Inches(10)
        slide_w = int(slide_h * ratio)
    else:
        # Near-square or standard
        slide_h = Inches(7.5)
        slide_w = int(slide_h * ratio)
    return slide_w, slide_h


# ─── Mode: IMAGE ───────────────────────────────────────────────────

def _build_image_slide(slide, page, mat, slide_w, slide_h):
    """Render the entire page as a single high-res image."""
    pix = page.get_pixmap(matrix=mat, alpha=False)
    img_stream = io.BytesIO(pix.tobytes("png"))
    slide.shapes.add_picture(img_stream, 0, 0, slide_w, slide_h)


# ─── Mode: HYBRID ─────────────────────────────────────────────────

def _build_hybrid_slide(slide, page, pdf_doc, mat, slide_w, slide_h, pw, ph):
    """
    Hybrid mode: pixel-perfect image + editable text hidden BEHIND it.
    User visually sees only the clean image (no ghost text), but the text
    is present in the slide — searchable (Ctrl+F), copyable, and editable
    after removing the image. This is how Adobe Acrobat "searchable PDF"
    layers work.
    """
    scale_x = slide_w / pw
    scale_y = slide_h / ph

    # 1. Add text FIRST so it sits BEHIND the image (z-order = insertion order)
    _add_text(slide, page, scale_x, scale_y, transparent=True)

    # 2. Full-page image ON TOP — covers the text visually
    pix = page.get_pixmap(matrix=mat, alpha=False)
    img_stream = io.BytesIO(pix.tobytes("png"))
    slide.shapes.add_picture(img_stream, 0, 0, slide_w, slide_h)


# ─── Mode: EDITABLE ───────────────────────────────────────────────

def _build_editable_slide(slide, page, pdf_doc, slide_w, slide_h, pw, ph):
    """Fully editable slide with native PPT elements."""
    scale_x = slide_w / pw
    scale_y = slide_h / ph

    # 1. White background rectangle
    bg = slide.shapes.add_shape(1, 0, 0, slide_w, slide_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    # 2. Colored rectangles and shapes (backgrounds, banners, dividers)
    _add_rects(slide, page, scale_x, scale_y)

    # 3. Lines and paths (borders, separators, decorative lines)
    _add_lines(slide, page, scale_x, scale_y)

    # 4. Images at correct positions
    _add_images(slide, page, pdf_doc, scale_x, scale_y)

    # 5. Editable text (on top of everything) — visible, real colors
    _add_text(slide, page, scale_x, scale_y, transparent=False)


# ─── Shape extraction ──────────────────────────────────────────────

def _add_rects(slide, page, scale_x, scale_y):
    """Add colored rectangles and filled shapes from PDF drawings."""
    page_w = page.rect.width
    page_h = page.rect.height

    seen = set()
    for path in page.get_drawings():
        fill = path.get("fill")
        rect = path.get("rect")
        if rect is None:
            continue

        # Dedup identical shapes (some PDFs double-draw)
        key = (round(rect.x0, 1), round(rect.y0, 1),
               round(rect.x1, 1), round(rect.y1, 1),
               tuple(round(c, 2) for c in (fill or path.get("color") or ())))
        if key in seen:
            continue
        seen.add(key)

        # Skip near-page-size backgrounds (they're just white canvas)
        if rect.width > page_w * 0.92 and rect.height > page_h * 0.92:
            continue

        # Filled shapes
        if fill is not None:
            # Skip near-white fills (backgrounds)
            if all(c > 0.95 for c in fill[:3]):
                continue
            if rect.width < 3 or rect.height < 3:
                continue

            left = int(rect.x0 * scale_x)
            top = int(rect.y0 * scale_y)
            width = int(rect.width * scale_x)
            height = int(rect.height * scale_y)

            try:
                shape = slide.shapes.add_shape(1, left, top, width, height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(
                    min(255, int(fill[0] * 255)),
                    min(255, int(fill[1] * 255)),
                    min(255, int(fill[2] * 255)),
                )
                shape.line.fill.background()
                shape.rotation = 0.0
            except Exception:
                pass
            continue

        # Stroked-only shapes (borders, frames)
        stroke = path.get("color")
        if stroke is None:
            continue
        if rect.width < 5 or rect.height < 5:
            continue

        left = int(rect.x0 * scale_x)
        top = int(rect.y0 * scale_y)
        width = int(rect.width * scale_x)
        height = int(rect.height * scale_y)

        try:
            shape = slide.shapes.add_shape(1, left, top, width, height)
            shape.fill.background()
            shape.line.color.rgb = RGBColor(
                min(255, int(stroke[0] * 255)),
                min(255, int(stroke[1] * 255)),
                min(255, int(stroke[2] * 255)),
            )
            lw = path.get("width", 1)
            shape.line.width = Pt(max(0.25, lw))
        except Exception:
            pass


def _add_lines(slide, page, scale_x, scale_y):
    """Add lines and connectors from PDF drawings (borders, separators, rules)."""
    seen = set()
    for path in page.get_drawings():
        items = path.get("items", [])
        color = path.get("color")
        if not items or color is None:
            continue
        # Skip if it has a fill (handled by _add_rects)
        if path.get("fill") is not None:
            continue

        # Emit every line segment in the path (not just single-segment paths)
        # so polyline borders and frame edges are preserved too.
        for item in items:
            if item[0] != "l":
                continue
            p1, p2 = item[1], item[2]

            # Skip degenerate lines shorter than 3pt in PDF
            dx_pdf = p2.x - p1.x
            dy_pdf = p2.y - p1.y
            if (dx_pdf * dx_pdf + dy_pdf * dy_pdf) < 9:
                continue

            x1 = int(p1.x * scale_x)
            y1 = int(p1.y * scale_y)
            x2 = int(p2.x * scale_x)
            y2 = int(p2.y * scale_y)

            # Deduplicate lines drawn multiple times in the source PDF
            key = (x1, y1, x2, y2)
            if key in seen:
                continue
            seen.add(key)

            try:
                connector = slide.shapes.add_connector(
                    1, x1, y1, x2 - x1, y2 - y1  # MSO_CONNECTOR_TYPE.STRAIGHT
                )
                connector.line.color.rgb = RGBColor(
                    min(255, int(color[0] * 255)),
                    min(255, int(color[1] * 255)),
                    min(255, int(color[2] * 255)),
                )
                lw = path.get("width", 1)
                connector.line.width = Pt(max(0.25, lw))
            except Exception:
                pass


# ─── Image extraction ──────────────────────────────────────────────

def _add_images(slide, page, pdf_doc, scale_x, scale_y):
    """Extract and place images at their exact PDF positions."""
    seen = set()
    for img_info in page.get_images(full=True):
        xref = img_info[0]
        if xref in seen:
            continue
        seen.add(xref)

        try:
            base = pdf_doc.extract_image(xref)
            if not base or "image" not in base:
                continue
            img_data = base["image"]
            img_ext = base.get("ext", "png")

            # Some PDFs have CMYK or unusual color spaces - convert via raw pixmap
            if img_ext in ("jpx", "jb2", "ccitt"):
                pix = fitz.Pixmap(pdf_doc, xref)
                if pix.n > 4:
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                img_data = pix.tobytes("png")

            for rect in page.get_image_rects(xref):
                if rect.width < 2 or rect.height < 2:
                    continue

                left = int(rect.x0 * scale_x)
                top = int(rect.y0 * scale_y)
                width = int(rect.width * scale_x)
                height = int(rect.height * scale_y)

                if width < Emu(10000) or height < Emu(10000):
                    continue

                slide.shapes.add_picture(
                    io.BytesIO(img_data), left, top, width, height
                )
        except Exception:
            pass


# ─── Text extraction ───────────────────────────────────────────────

# Common font name mappings for better PowerPoint compatibility
_FONT_MAP = {
    "ArialMT": "Arial",
    "TimesNewRomanPSMT": "Times New Roman",
    "TimesNewRomanPS": "Times New Roman",
    "CourierNewPSMT": "Courier New",
    "Helvetica": "Arial",
    "HelveticaNeue": "Arial",
    "Verdana": "Verdana",
    "Georgia": "Georgia",
    "Calibri": "Calibri",
    "Cambria": "Cambria",
    "Tahoma": "Tahoma",
    "TrebuchetMS": "Trebuchet MS",
    "Symbol": "Symbol",
    "ZapfDingbats": "Wingdings",
}


def _clean_font_name(raw_font):
    """Clean and map PDF font name to a PowerPoint-compatible font."""
    if not raw_font:
        return "Arial"

    # Remove subset prefix (e.g. ABCDEF+FontName)
    if "+" in raw_font:
        raw_font = raw_font.split("+", 1)[1]

    # Check direct mapping first
    base_key = raw_font.split("-")[0].split(",")[0].strip()
    if base_key in _FONT_MAP:
        return _FONT_MAP[base_key]

    # Strip style suffixes
    clean = base_key
    for suffix in ("Bold", "Italic", "Oblique", "Regular", "Medium",
                    "Light", "Semibold", "SemiBold", "ExtraBold", "Black",
                    "Thin", "Book", "Roman", "Condensed", "Narrow", "MT", "PS"):
        clean = clean.replace(suffix, "")
    clean = clean.strip()

    if not clean:
        return "Arial"

    # Add spaces before capitals (e.g. TimesNewRoman -> Times New Roman)
    spaced = ""
    for i, ch in enumerate(clean):
        if i > 0 and ch.isupper() and clean[i - 1].islower():
            spaced += " "
        spaced += ch

    return spaced.strip() or "Arial"


def _detect_alignment(spans, block_bbox):
    """Detect text alignment based on span positions within the block."""
    if not spans:
        return PP_ALIGN.LEFT

    block_width = block_bbox[2] - block_bbox[0]
    if block_width < 10:
        return PP_ALIGN.LEFT

    # Check first span position relative to block
    first_span = spans[0]
    span_bbox = first_span.get("bbox", block_bbox)
    left_margin = span_bbox[0] - block_bbox[0]
    right_margin = block_bbox[2] - span_bbox[2]

    # Centered text: roughly equal margins on both sides
    if block_width > 50 and abs(left_margin - right_margin) < block_width * 0.15:
        if left_margin > block_width * 0.1:
            return PP_ALIGN.CENTER

    # Right-aligned: large left margin, small right margin
    if left_margin > block_width * 0.4 and right_margin < block_width * 0.1:
        return PP_ALIGN.RIGHT

    return PP_ALIGN.LEFT


def _add_text(slide, page, scale_x, scale_y, transparent=True):
    """
    Line-precise text placement: one textbox per visible line.

    This avoids block-level stacking problems where multiple columns of text
    on the same page (common in tickets, receipts, multi-column layouts) get
    lumped together into a single textbox and end up overlapping.
    """
    text_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
    links = _rects_for_links(page)

    for block in text_dict.get("blocks", []):
        if block.get("type") != 0:
            continue

        for line in block.get("lines", []):
            spans = [s for s in line.get("spans", []) if s.get("text", "")]
            if not spans:
                continue
            if not any(s.get("text", "").strip() for s in spans):
                continue

            line_bbox = line.get("bbox", [0, 0, 0, 0])
            x0, y0, x1, y1 = line_bbox
            pdf_w = x1 - x0
            pdf_h = y1 - y0
            if pdf_w < 0.5 or pdf_h < 0.5:
                continue

            # Tight fit: horizontal padding to prevent clipping, a tiny
            # vertical expansion so the top of ascenders isn't clipped.
            left = int(x0 * scale_x)
            top = int((y0 - 0.5) * scale_y)
            width = int(pdf_w * scale_x * 1.08 + Emu(50000))
            height = int((pdf_h + 1.0) * scale_y * 1.25)

            if top < 0:
                top = 0

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = False  # per-line textbox; don't wrap
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            tf.auto_size = MSO_AUTO_SIZE.NONE  # keep fixed size — no auto-growth
            tf.vertical_anchor = MSO_ANCHOR.TOP

            # Transparent textbox fill & no border
            txBox.fill.background()
            txBox.line.fill.background()

            alignment = _detect_alignment(spans, line_bbox)

            p = tf.paragraphs[0]
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.0
            p.alignment = alignment

            for span_idx, span in enumerate(spans):
                text = span.get("text", "")
                if not text:
                    continue

                # Detect & convert leading bullet glyph on the first span
                # of a line — PowerPoint renders bullets via paragraph
                # formatting, but this keeps native bullet chars readable.
                if span_idx == 0:
                    bullet, text = _strip_bullet_prefix(text)
                    if bullet:
                        prefix_run = p.add_run()
                        prefix_run.text = f"{bullet}  "
                        size = span.get("size", 12)
                        prefix_run.font.size = Pt(max(1, size))
                        prefix_run.font.name = "Arial"
                        if transparent:
                            _set_run_invisible(prefix_run)

                if not text:
                    continue

                run = p.add_run()
                run.text = text

                size = span.get("size", 12)

                # Superscript / subscript detection from flag bits
                flags = span.get("flags", 0)
                is_super = bool(flags & 1)         # bit 0
                is_sub   = bool(flags & 8) and not is_super   # uncommon — some PDFs

                if is_super or is_sub:
                    # PowerPoint represents super/sub via baseline offset
                    run.font.size = Pt(max(1, size * 0.7))
                else:
                    run.font.size = Pt(max(1, size))

                font_name = span.get("font", "")
                run.font.name = _clean_font_name(font_name)

                font_lower = font_name.lower()
                if "bold" in font_lower or "black" in font_lower or (flags & 16):
                    run.font.bold = True
                if "italic" in font_lower or "oblique" in font_lower or (flags & 2):
                    run.font.italic = True
                if flags & 4:
                    run.font.underline = True

                color_int = span.get("color", 0)
                if color_int:
                    r_val = (color_int >> 16) & 0xFF
                    g_val = (color_int >> 8) & 0xFF
                    b_val = color_int & 0xFF
                    run.font.color.rgb = RGBColor(r_val, g_val, b_val)
                else:
                    run.font.color.rgb = RGBColor(0, 0, 0)

                # Hyperlink: if this span sits inside a URI link annotation,
                # make it clickable in PowerPoint.
                span_bbox = span.get("bbox")
                if span_bbox and links:
                    uri = _find_link_for_span(span_bbox, links)
                    if uri:
                        _set_run_hyperlink(run, uri)

                # Apply super/sub baseline offset via OOXML (python-pptx
                # doesn't expose this property directly).
                if is_super or is_sub:
                    try:
                        rPr = run._r.get_or_add_rPr()
                        rPr.set("baseline", "30000" if is_super else "-25000")
                    except Exception:
                        pass

                # In hybrid mode: make run fully transparent (alpha=0). Text
                # stays selectable / searchable. The image on top provides
                # the visual — this matches Adobe Acrobat's OCR layer pattern.
                if transparent:
                    _set_run_invisible(run)
