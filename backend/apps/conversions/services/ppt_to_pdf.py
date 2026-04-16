import io
import fitz
from pptx import Presentation
from pptx.util import Emu


def convert_ppt_to_pdf(file_path):
    """
    Convert PowerPoint presentation to PDF.
    Each slide becomes a page.
    """
    prs = Presentation(file_path)
    pdf_doc = fitz.open()

    slide_width_pt = prs.slide_width / 914400 * 72  # EMU to points
    slide_height_pt = prs.slide_height / 914400 * 72

    for slide in prs.slides:
        page = pdf_doc.new_page(width=slide_width_pt, height=slide_height_pt)

        for shape in slide.shapes:
            # Position
            left = (shape.left or 0) / 914400 * 72
            top = (shape.top or 0) / 914400 * 72

            if shape.has_text_frame:
                y_offset = top
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        y_offset += 14
                        continue

                    font_size = 12
                    font_name = "helv"
                    is_bold = False

                    if para.runs:
                        run = para.runs[0]
                        if run.font.size:
                            font_size = run.font.size.pt
                        if run.font.bold:
                            is_bold = True
                            font_name = "helvB"

                    page.insert_text(
                        fitz.Point(left, y_offset + font_size),
                        text,
                        fontsize=font_size,
                        fontname=font_name,
                    )
                    y_offset += font_size * 1.4

            elif hasattr(shape, "image"):
                try:
                    img_bytes = shape.image.blob
                    width = (shape.width or Emu(1)) / 914400 * 72
                    height = (shape.height or Emu(1)) / 914400 * 72
                    rect = fitz.Rect(left, top, left + width, top + height)
                    page.insert_image(rect, stream=img_bytes)
                except Exception:
                    pass

    output = io.BytesIO()
    pdf_doc.save(output)
    pdf_doc.close()
    return output.getvalue()
